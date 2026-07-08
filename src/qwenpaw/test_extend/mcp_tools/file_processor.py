# -*- coding: utf-8 -*-
"""File type detection and content extraction for agent file uploads.

Provides tools to handle non-image file uploads (documents, spreadsheets, etc.)
by extracting text content before sending to the model.
"""

import os
from pathlib import Path
from typing import Union


def detect_file_type(filename: str) -> str:
    """Detect file type category from filename extension.

    Returns one of: image, document, spreadsheet, text, code, pdf, unknown
    """
    ext = Path(filename).suffix.lower()
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg"}:
        return "image"
    if ext in {".md", ".markdown", ".txt", ".text"}:
        return "text"
    if ext in {".csv", ".tsv"}:
        return "spreadsheet"
    if ext in {".xlsx", ".xls"}:
        return "spreadsheet"
    if ext in {".docx", ".doc"}:
        return "document"
    if ext in {".pptx", ".ppt"}:
        return "presentation"
    if ext in {".pdf"}:
        return "pdf"
    if ext in {".json", ".yaml", ".yml", ".xml", ".toml", ".ini", ".cfg"}:
        return "text"
    if ext in {".py", ".js", ".ts", ".java", ".go", ".rs", ".c", ".cpp", ".h"}:
        return "code"
    return "unknown"


async def process_uploaded_file_tool(file_path: str, filename: str = "") -> dict:
    """Process an uploaded file and extract text content for the model.

    For images: returns metadata only (model can read directly).
    For documents/spreadsheets/text: extracts and returns text content.

    Args:
        file_path: Absolute path to the uploaded file.
        filename: Original filename (for type detection).

    Returns:
        Dict with file_type, content (extracted text or None for images),
        and a message describing what was extracted.
    """
    if not filename:
        filename = os.path.basename(file_path)

    if not os.path.exists(file_path):
        return {
            "file_type": "error",
            "content": None,
            "message": f"文件不存在: {file_path}",
        }

    file_type = detect_file_type(filename)
    file_size = os.path.getsize(file_path)

    result = {
        "file_type": file_type,
        "filename": filename,
        "file_size": file_size,
        "content": None,
        "message": "",
    }

    if file_type == "image":
        result["message"] = f"图片文件 ({filename}, {file_size} 字节)，可直接发送给模型"
        return result

    if file_type == "text":
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            result["content"] = content[:50000]  # Limit to 50KB
            result["message"] = f"文本文件 ({filename})，已提取 {len(content)} 字符"
        except Exception as e:
            result["message"] = f"读取失败: {e}"
        return result

    if file_type == "spreadsheet":
        ext = Path(filename).suffix.lower()
        if ext == ".csv" or ext == ".tsv":
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
                result["content"] = content[:50000]
                result["message"] = f"CSV 文件 ({filename})，已提取 {len(content)} 字符"
            except Exception as e:
                result["message"] = f"读取失败: {e}"
        else:
            # xlsx / xls
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        row_str = "\t".join(str(c) if c is not None else "" for c in row)
                        rows.append(row_str)
                    parts.append(f"Sheet: {sheet_name}\n" + "\n".join(rows))
                wb.close()
                content = "\n\n".join(parts)
                result["content"] = content[:50000]
                result["message"] = f"Excel 文件 ({filename})，已提取 {len(parts)} 个工作表"
            except ImportError:
                result["message"] = f"Excel 文件 ({filename})，缺少 openpyxl 库无法解析"
            except Exception as e:
                result["message"] = f"Excel 解析失败: {e}"
        return result

    if file_type == "document":
        ext = Path(filename).suffix.lower()
        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs]
                content = "\n".join(paragraphs)
                result["content"] = content[:50000]
                result["message"] = f"Word 文档 ({filename})，已提取 {len(paragraphs)} 段落"
            except ImportError:
                result["message"] = f"Word 文档 ({filename})，缺少 python-docx 库无法解析"
            except Exception as e:
                result["message"] = f"Word 解析失败: {e}"
        else:
            # .doc (older format)
            result["message"] = f"旧版 .doc 格式 ({filename})，建议转换为 .docx"
        return result

    if file_type == "pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            parts = []
            for i, page in enumerate(doc):
                text = page.get_text()
                parts.append(f"--- 第 {i+1} 页 ---\n{text}")
            doc.close()
            content = "\n\n".join(parts)
            result["content"] = content[:50000]
            result["message"] = f"PDF 文件 ({filename})，已提取 {len(parts)} 页"
        except ImportError:
            result["message"] = f"PDF 文件 ({filename})，缺少 PyMuPDF 库无法解析"
        except Exception as e:
            result["message"] = f"PDF 解析失败: {e}"
        return result

    if file_type == "presentation":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
                parts.append(f"Slide {i+1}:\n" + "\n".join(texts))
            content = "\n\n".join(parts)
            result["content"] = content[:50000]
            result["message"] = f"PPT 文件 ({filename})，已提取 {len(parts)} 页"
        except ImportError:
            result["message"] = f"PPT 文件 ({filename})，缺少 python-pptx 库无法解析"
        except Exception as e:
            result["message"] = f"PPT 解析失败: {e}"
        return result

    result["message"] = f"未知文件类型 ({filename})，无法自动提取内容"
    return result
